"""
Web Fetch tool - Fetch and extract readable content from web pages and remote files.

Supports:
- HTML web pages: extracts readable text content
- Document files (PDF, Word, TXT, Markdown, etc.): downloads to workspace/tmp and parses content
"""

import html as html_lib
import os
import re
import uuid
from typing import Dict, Any, Optional, Set
from urllib.parse import urlparse, unquote, urljoin

import requests

from agent.tools.base_tool import BaseTool, ToolResult
from agent.tools.utils.truncate import truncate_head, format_size
from agent.tools.utils.url_safety import validate_url_safe
from common.log import logger


DEFAULT_TIMEOUT = 30
MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB
# Cap on how many redirects we follow; each hop's target is re-validated
# against the SSRF guard so a public URL cannot bounce us into an internal one.
MAX_REDIRECTS = 10

DEFAULT_HEADERS = {
    "User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36",
    "Accept": "*/*",
}

# Supported document file extensions
PDF_SUFFIXES: Set[str] = {".pdf"}
WORD_SUFFIXES: Set[str] = {".docx"}
TEXT_SUFFIXES: Set[str] = {".txt", ".md", ".markdown", ".rst", ".csv", ".tsv", ".log"}
SPREADSHEET_SUFFIXES: Set[str] = {".xls", ".xlsx"}
PPT_SUFFIXES: Set[str] = {".ppt", ".pptx"}

ALL_DOC_SUFFIXES = PDF_SUFFIXES | WORD_SUFFIXES | TEXT_SUFFIXES | SPREADSHEET_SUFFIXES | PPT_SUFFIXES
RECOVERY_LINK_LIMIT = 16
RECOVERY_PAGE_LIMIT = 3

_CHARSET_RE = re.compile(r'charset\s*=\s*["\']?\s*([\w\-]+)', re.IGNORECASE)
_META_CHARSET_RE = re.compile(rb'<meta[^>]+charset\s*=\s*["\']?\s*([\w\-]+)', re.IGNORECASE)
_META_HTTP_EQUIV_RE = re.compile(
    rb'<meta[^>]+http-equiv\s*=\s*["\']?Content-Type["\']?[^>]+content\s*=\s*["\'][^"\']*charset=([\w\-]+)',
    re.IGNORECASE,
)


def _extract_charset_from_content_type(content_type: str) -> Optional[str]:
    """Extract charset from Content-Type header value."""
    m = _CHARSET_RE.search(content_type)
    return m.group(1) if m else None


def _extract_charset_from_html_meta(raw_bytes: bytes) -> Optional[str]:
    """Extract charset from HTML <meta> tags in the first few KB of raw bytes."""
    m = _META_CHARSET_RE.search(raw_bytes)
    if m:
        return m.group(1).decode("ascii", errors="ignore")
    m = _META_HTTP_EQUIV_RE.search(raw_bytes)
    if m:
        return m.group(1).decode("ascii", errors="ignore")
    return None


def _get_url_suffix(url: str) -> str:
    """Extract file extension from URL path, ignoring query params."""
    path = urlparse(url).path
    return os.path.splitext(path)[-1].lower()


def _is_document_url(url: str) -> bool:
    """Check if URL points to a downloadable document file."""
    suffix = _get_url_suffix(url)
    return suffix in ALL_DOC_SUFFIXES


class WebFetch(BaseTool):
    """Tool for fetching web pages and remote document files"""

    name: str = "web_fetch"
    description: str = (
        "Fetch content from a http/https URL. For web pages, extracts readable text. "
        "For document files (PDF, Word, TXT, Markdown, Excel, PPT), downloads and parses the file content. "
        "Supported file types: .pdf, .docx, .txt, .md, .csv, .xls, .xlsx, .ppt, .pptx"
    )

    params: dict = {
        "type": "object",
        "properties": {
            "url": {
                "type": "string",
                "description": "The HTTP/HTTPS URL to fetch (web page or document file link)"
            }
        },
        "required": ["url"]
    }

    def __init__(self, config: dict = None):
        self.config = config or {}
        self.cwd = self.config.get("cwd", os.getcwd())

    def execute(self, args: Dict[str, Any]) -> ToolResult:
        url = args.get("url", "").strip()
        if not url:
            return ToolResult.fail("Error: 'url' parameter is required")

        parsed = urlparse(url)
        if parsed.scheme not in ("http", "https"):
            return ToolResult.fail("Error: Invalid URL (must start with http:// or https://)")

        # SSRF guard: reject URLs that resolve to private/loopback/link-local/
        # cloud-metadata addresses before any request is issued.
        try:
            validate_url_safe(url)
        except ValueError as e:
            return ToolResult.fail(f"Error: {e}")

        if _is_document_url(url):
            return self._fetch_document(url)

        return self._fetch_webpage(url)

    # ---- Safe request helper ----

    def _configured_proxy(self) -> Optional[str]:
        """Resolve proxy for web_fetch requests.

        Priority:
        1. tools.web_fetch.proxy
        2. global proxy
        3. requests' default environment proxy handling
        """
        if isinstance(self.config, dict):
            proxy = (self.config.get("proxy") or "").strip()
            if proxy:
                return proxy

        try:
            from config import conf
            proxy = (conf().get("proxy") or "").strip()
            return proxy or None
        except Exception:
            return None

    def _request_proxies(self) -> Optional[Dict[str, str]]:
        proxy = self._configured_proxy()
        return {"http": proxy, "https": proxy} if proxy else None

    def _safe_get(self, url: str, **kwargs) -> requests.Response:
        """Issue a GET request while re-validating every redirect hop (SSRF guard).

        Auto-redirect is disabled and each hop is followed manually so the
        target of every redirect is re-resolved and checked against the SSRF
        guard. This prevents a public URL from 3xx-bouncing into a private,
        loopback, link-local or cloud-metadata address. ``kwargs`` are passed
        through to ``requests.get`` (e.g. ``stream``).

        Raises:
            ValueError: if any hop resolves to a non-public address.
        """
        kwargs.pop("allow_redirects", None)
        current = url
        proxies = self._request_proxies()
        for _ in range(MAX_REDIRECTS + 1):
            request_kwargs = dict(kwargs)
            if proxies:
                request_kwargs["proxies"] = proxies
            response = requests.get(
                current,
                headers=DEFAULT_HEADERS,
                timeout=DEFAULT_TIMEOUT,
                allow_redirects=False,
                **request_kwargs,
            )
            if not response.is_redirect and not response.is_permanent_redirect:
                return response

            location = response.headers.get("Location")
            if not location:
                return response

            # Resolve the redirect target relative to the current URL, then
            # re-validate it before following.
            current = requests.compat.urljoin(current, location)
            validate_url_safe(current)
            response.close()

        raise ValueError(f"Too many redirects (>{MAX_REDIRECTS})")

    # ---- Web page fetching ----

    def _fetch_webpage(self, url: str) -> ToolResult:
        """Fetch and extract readable text from an HTML web page."""
        parsed = urlparse(url)
        try:
            response = self._safe_get(url)
            response.raise_for_status()
        except requests.Timeout:
            return ToolResult.fail(f"Error: Request timed out after {DEFAULT_TIMEOUT}s")
        except requests.ConnectionError:
            return ToolResult.fail(f"Error: Failed to connect to {parsed.netloc}")
        except requests.HTTPError as e:
            status_code = e.response.status_code if e.response is not None else "unknown"
            hint = self._build_http_error_hint(url, status_code)
            return ToolResult.fail(f"Error: HTTP {status_code} for URL: {url}{hint}")
        except ValueError as e:
            return ToolResult.fail(f"Error: {e}")
        except Exception as e:
            return ToolResult.fail(f"Error: Failed to fetch URL: {e}")

        content_type = response.headers.get("Content-Type", "")
        if self._is_binary_content_type(content_type) and not _is_document_url(url):
            return self._handle_download_by_content_type(url, response, content_type)

        response.encoding = self._detect_encoding(response)
        html = response.text
        title = self._extract_title(html)
        text = self._extract_text(html)

        return ToolResult.success(f"Title: {title}\n\nContent:\n{text}")

    # ---- Document fetching ----

    def _fetch_document(self, url: str) -> ToolResult:
        """Download a document file and extract its text content."""
        suffix = _get_url_suffix(url)
        parsed = urlparse(url)
        filename = self._extract_filename(url)
        tmp_dir = self._ensure_tmp_dir()

        local_path = os.path.join(tmp_dir, filename)
        logger.info(f"[WebFetch] Downloading document: {url} -> {local_path}")

        try:
            response = self._safe_get(url, stream=True)
            response.raise_for_status()

            content_length = int(response.headers.get("Content-Length", 0))
            if content_length > MAX_FILE_SIZE:
                return ToolResult.fail(
                    f"Error: File too large ({format_size(content_length)} > {format_size(MAX_FILE_SIZE)})"
                )

            downloaded = 0
            with open(local_path, "wb") as f:
                for chunk in response.iter_content(chunk_size=8192):
                    downloaded += len(chunk)
                    if downloaded > MAX_FILE_SIZE:
                        f.close()
                        os.remove(local_path)
                        return ToolResult.fail(
                            f"Error: File too large (>{format_size(MAX_FILE_SIZE)}), download aborted"
                        )
                    f.write(chunk)

        except requests.Timeout:
            return ToolResult.fail(f"Error: Download timed out after {DEFAULT_TIMEOUT}s")
        except requests.ConnectionError:
            return ToolResult.fail(f"Error: Failed to connect to {parsed.netloc}")
        except requests.HTTPError as e:
            return ToolResult.fail(f"Error: HTTP {e.response.status_code} for URL: {url}")
        except ValueError as e:
            self._cleanup_file(local_path)
            return ToolResult.fail(f"Error: {e}")
        except Exception as e:
            self._cleanup_file(local_path)
            return ToolResult.fail(f"Error: Failed to download file: {e}")

        try:
            text = self._parse_document(local_path, suffix)
        except Exception as e:
            self._cleanup_file(local_path)
            return ToolResult.fail(f"Error: Failed to parse document: {e}")

        if not text or not text.strip():
            file_size = os.path.getsize(local_path)
            return ToolResult.success(
                f"File downloaded to: {local_path} ({format_size(file_size)})\n"
                f"No text content could be extracted. The file may contain only images or be encrypted."
            )

        truncation = truncate_head(text)
        result_text = truncation.content

        file_size = os.path.getsize(local_path)
        header = f"[Document: {filename} | Size: {format_size(file_size)} | Saved to: {local_path}]\n\n"

        if truncation.truncated:
            header += f"[Content truncated: showing {truncation.output_lines} of {truncation.total_lines} lines]\n\n"

        return ToolResult.success(header + result_text)

    def _parse_document(self, file_path: str, suffix: str) -> str:
        """Parse document file and return extracted text."""
        if suffix in PDF_SUFFIXES:
            return self._parse_pdf(file_path)
        elif suffix in WORD_SUFFIXES:
            return self._parse_word(file_path)
        elif suffix in TEXT_SUFFIXES:
            return self._parse_text(file_path)
        elif suffix in SPREADSHEET_SUFFIXES:
            return self._parse_spreadsheet(file_path)
        elif suffix in PPT_SUFFIXES:
            return self._parse_ppt(file_path)
        else:
            return self._parse_text(file_path)

    def _parse_pdf(self, file_path: str) -> str:
        """Extract text from PDF using pypdf."""
        try:
            from pypdf import PdfReader
        except ImportError:
            raise ImportError("pypdf library is required for PDF parsing. Install with: pip install pypdf")

        reader = PdfReader(file_path)
        text_parts = []
        for page_num, page in enumerate(reader.pages, 1):
            page_text = page.extract_text()
            if page_text and page_text.strip():
                text_parts.append(f"--- Page {page_num}/{len(reader.pages)} ---\n{page_text}")

        return "\n\n".join(text_parts)

    def _parse_word(self, file_path: str) -> str:
        """Extract text from Word documents (.docx)."""
        try:
            from docx import Document
        except ImportError:
            raise ImportError(
                "python-docx library is required for .docx parsing. Install with: pip install python-docx"
            )
        doc = Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n\n".join(paragraphs)

    def _parse_text(self, file_path: str) -> str:
        """Read plain text files (txt, md, csv, etc.)."""
        encodings = ["utf-8", "utf-8-sig", "gbk", "gb2312", "latin-1"]
        for enc in encodings:
            try:
                with open(file_path, "r", encoding=enc) as f:
                    return f.read()
            except (UnicodeDecodeError, UnicodeError):
                continue
        raise ValueError(f"Unable to decode file with any supported encoding: {encodings}")

    def _parse_spreadsheet(self, file_path: str) -> str:
        """Extract text from Excel files (.xls/.xlsx)."""
        try:
            import openpyxl
        except ImportError:
            raise ImportError(
                "openpyxl library is required for .xlsx parsing. Install with: pip install openpyxl"
            )

        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        result_parts = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    rows.append(" | ".join(cells))
            if rows:
                result_parts.append(f"--- Sheet: {sheet_name} ---\n" + "\n".join(rows))

        wb.close()
        return "\n\n".join(result_parts)

    def _parse_ppt(self, file_path: str) -> str:
        """Extract text from PowerPoint files (.ppt/.pptx)."""
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError(
                "python-pptx library is required for .pptx parsing. Install with: pip install python-pptx"
            )

        prs = Presentation(file_path)
        text_parts = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_texts.append(text)
            if slide_texts:
                text_parts.append(f"--- Slide {slide_num}/{len(prs.slides)} ---\n" + "\n".join(slide_texts))

        return "\n\n".join(text_parts)

    # ---- Encoding detection ----

    @staticmethod
    def _detect_encoding(response: requests.Response) -> str:
        """Detect response encoding with priority: Content-Type header > HTML meta > chardet > utf-8."""
        # 1. Check Content-Type header for explicit charset
        content_type = response.headers.get("Content-Type", "")
        charset = _extract_charset_from_content_type(content_type)
        if charset:
            return charset

        # 2. Scan raw bytes for HTML meta charset declaration
        raw = response.content[:4096]
        charset = _extract_charset_from_html_meta(raw)
        if charset:
            return charset

        # 3. Use apparent_encoding (chardet-based detection) if confident enough
        apparent = response.apparent_encoding
        if apparent:
            apparent_lower = apparent.lower()
            # Trust CJK / Windows encodings detected by chardet
            trusted_prefixes = ("utf", "gb", "big5", "euc", "shift_jis", "iso-2022", "windows", "ascii")
            if any(apparent_lower.startswith(p) for p in trusted_prefixes):
                return apparent

        # 4. Fallback
        return "utf-8"

    # ---- Helper methods ----

    def _build_http_error_hint(self, url: str, status_code) -> str:
        """Return recovery guidance and real same-site links for missing pages."""
        if status_code not in (404, 410):
            return ""

        pages = []
        for fallback_url in self._candidate_fallback_urls(url):
            try:
                response = self._safe_get(fallback_url)
                response.raise_for_status()
            except Exception as e:
                logger.debug(f"[WebFetch] Recovery fetch failed: {fallback_url} ({e})")
                continue

            content_type = response.headers.get("Content-Type", "")
            if self._is_binary_content_type(content_type):
                continue

            response.encoding = self._detect_encoding(response)
            html = response.text
            links = self._extract_links(html, fallback_url, RECOVERY_LINK_LIMIT)
            if links:
                pages.append((fallback_url, self._extract_title(html), links))
            if len(pages) >= RECOVERY_PAGE_LIMIT:
                break

        guidance = [
            "",
            "",
            "Recovery hint: this page does not exist. Do not keep guessing URL paths.",
            "Open one of the discovered same-site navigation links below, or summarize content already fetched and ask the user for a more specific link.",
        ]

        if pages:
            guidance.append("Discovered navigation links:")
            seen_links = set()
            for page_url, title, links in pages:
                guidance.append(f"- From {title or 'Untitled'} ({page_url}):")
                for text, href in links:
                    if href in seen_links:
                        continue
                    seen_links.add(href)
                    guidance.append(f"  - {text}: {href}")
        else:
            guidance.append("No fallback navigation links could be extracted from parent pages.")

        return "\n" + "\n".join(guidance)

    @staticmethod
    def _candidate_fallback_urls(url: str):
        parsed = urlparse(url)
        base = f"{parsed.scheme}://{parsed.netloc}"
        path = parsed.path or "/"
        if not path.startswith("/"):
            path = "/" + path

        candidates = []
        seen = set()
        parts = [part for part in path.strip("/").split("/") if part]
        for count in range(len(parts) - 1, -1, -1):
            parent_path = "/" + "/".join(parts[:count])
            if parent_path != "/":
                parent_path += "/"
            candidate = base + parent_path
            if candidate == url or candidate in seen:
                continue
            seen.add(candidate)
            candidates.append(candidate)
            if len(candidates) >= RECOVERY_PAGE_LIMIT:
                break

        return candidates

    @staticmethod
    def _extract_links(html: str, base_url: str, limit: int = RECOVERY_LINK_LIMIT):
        base_parsed = urlparse(base_url)
        links = []
        seen = set()
        anchor_re = re.compile(
            r"<a\b[^>]*\bhref\s*=\s*(['\"])(.*?)\1[^>]*>(.*?)</a>",
            re.IGNORECASE | re.DOTALL,
        )
        for match in anchor_re.finditer(html):
            href = html_lib.unescape(match.group(2).strip())
            if not href or href.startswith(("#", "javascript:", "mailto:", "tel:")):
                continue

            absolute = urljoin(base_url, href)
            parsed = urlparse(absolute)
            if parsed.scheme not in ("http", "https") or parsed.netloc != base_parsed.netloc:
                continue

            try:
                validate_url_safe(absolute)
            except ValueError:
                continue

            normalized = parsed._replace(fragment="").geturl()
            if normalized in seen:
                continue

            text = re.sub(r"<[^>]+>", " ", match.group(3))
            text = html_lib.unescape(re.sub(r"\s+", " ", text)).strip()
            if not text:
                text = normalized

            seen.add(normalized)
            links.append((text[:80], normalized))
            if len(links) >= limit:
                break

        return links

    def _ensure_tmp_dir(self) -> str:
        """Ensure workspace/tmp directory exists and return its path."""
        tmp_dir = os.path.join(self.cwd, "tmp")
        os.makedirs(tmp_dir, exist_ok=True)
        return tmp_dir

    def _extract_filename(self, url: str) -> str:
        """Extract a safe filename from URL, with a short UUID prefix to avoid collisions."""
        path = urlparse(url).path
        basename = os.path.basename(unquote(path))
        if not basename or basename == "/":
            basename = "downloaded_file"
        # Sanitize: keep only safe chars
        basename = re.sub(r'[^\w.\-]', '_', basename)
        short_id = uuid.uuid4().hex[:8]
        return f"{short_id}_{basename}"

    @staticmethod
    def _cleanup_file(path: str):
        """Remove a file if it exists, ignoring errors."""
        try:
            if os.path.exists(path):
                os.remove(path)
        except Exception:
            pass

    @staticmethod
    def _is_binary_content_type(content_type: str) -> bool:
        """Check if Content-Type indicates a binary/document response."""
        binary_types = [
            "application/pdf",
            "application/vnd.openxmlformats",
            "application/vnd.ms-excel",
            "application/vnd.ms-powerpoint",
            "application/octet-stream",
        ]
        ct_lower = content_type.lower()
        return any(bt in ct_lower for bt in binary_types)

    def _handle_download_by_content_type(self, url: str, response: requests.Response, content_type: str) -> ToolResult:
        """Handle a URL that returned binary content instead of HTML."""
        ct_lower = content_type.lower()
        suffix_map = {
            "application/pdf": ".pdf",
            "application/vnd.openxmlformats-officedocument.wordprocessingml": ".docx",
            "application/vnd.ms-excel": ".xls",
            "application/vnd.openxmlformats-officedocument.spreadsheetml": ".xlsx",
            "application/vnd.ms-powerpoint": ".ppt",
            "application/vnd.openxmlformats-officedocument.presentationml": ".pptx",
        }
        detected_suffix = None
        for ct_prefix, ext in suffix_map.items():
            if ct_prefix in ct_lower:
                detected_suffix = ext
                break

        if detected_suffix and detected_suffix in ALL_DOC_SUFFIXES:
            # Re-fetch as document
            return self._fetch_document(url if _get_url_suffix(url) in ALL_DOC_SUFFIXES
                                        else self._rewrite_url_with_suffix(url, detected_suffix))
        return ToolResult.fail(f"Error: URL returned binary content ({content_type}), not a supported document type")

    @staticmethod
    def _rewrite_url_with_suffix(url: str, suffix: str) -> str:
        """Append a suffix to the URL path so _get_url_suffix works correctly."""
        parsed = urlparse(url)
        new_path = parsed.path.rstrip("/") + suffix
        return parsed._replace(path=new_path).geturl()

    # ---- HTML extraction (unchanged) ----

    @staticmethod
    def _extract_title(html: str) -> str:
        match = re.search(r"<title[^>]*>(.*?)</title>", html, re.IGNORECASE | re.DOTALL)
        return match.group(1).strip() if match else "Untitled"

    @staticmethod
    def _extract_text(html: str) -> str:
        text = re.sub(r"<script[^>]*>.*?</script>", "", html, flags=re.IGNORECASE | re.DOTALL)
        text = re.sub(r"<style[^>]*>.*?</style>", "", text, flags=re.IGNORECASE | re.DOTALL)
        text = re.sub(r"<[^>]+>", "", text)
        text = text.replace("&amp;", "&").replace("&lt;", "<").replace("&gt;", ">")
        text = text.replace("&quot;", '"').replace("&#39;", "'").replace("&nbsp;", " ")
        text = re.sub(r"[^\S\n]+", " ", text)
        text = re.sub(r"\n{3,}", "\n\n", text)
        lines = [line.strip() for line in text.splitlines()]
        text = "\n".join(lines)
        return text.strip()
